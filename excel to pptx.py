
import pandas as pd 
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
import wikipedia as wp
import urllib.request

# Download our data into a pandas dataframe 
# UPDATE ADDRESS TO WHERE YOUR CSV FILE IS 
df = pd.read_csv('E:/Code Project/Python - PPTX project/Nat_Park.csv')
del(df["Image"])
df['Location'] = df['Location'].map(str)
print(df.columns)

# Splitting data types into respective catagories 
df["Cordinates"] = df['Location'].str.split("/").str[1]
df["State"] = df['Location'].str.split('\r').str[0] 
df["State"] = df['State'].str.split('\n').str[0] 
df["Acres"] = df['Area (2021)'].str.split('acres').str[0]
df["KM"] = df['Area (2021)'].str.split('(').str[1]
df["KM"] = df["KM"].str.split('km').str[0]
df["Name"] = df["Name"].str.split('*').str[0]

# We can delete these unnecessary columns now 
del(df["Location"],df["Area (2021)"]) 



# function to convert strings into a list 
def Convert(string):
    li = list(string.split("src"))
    return li

# scraping all the photos off our wikipeida page so we can put them
# Using Beautiful soup to get all the possible photos on the wikipedia page 
import requests
URL = "https://en.wikipedia.org/wiki/List_of_national_parks_of_the_United_States#cite_note-:0-11" # Replace this with the website's URL
getURL = requests.get(URL, headers={"User-Agent":"Mozilla/5.0"})

from bs4 import BeautifulSoup

soup = BeautifulSoup(getURL.text, 'html.parser')

# converts all the 
images = soup.find_all('img')
# removed a bunch of the images url right before the table images on wikipedia 
images = str(images)[1870:]
ls = Convert(images)

# look for my images and split them up into a list 
ls = [k for k in ls if k.startswith("=\"//upload.wikimedia.org/wikipedia/commons/thu")]

# put the list into a pandas data frame for eaiser manipulation 
df2 = pd.DataFrame (ls, columns = ["url"])
# remove the bottom 63 since those url's are not important either 
df2 = df2.iloc[0:63:,] 
df2 = df2["url"].str[2:-2]

# add http: because we need it
df2 = "http:" +df2



# we can start using the python pptx library to create the powerpoints 

# create the powerpoint we want 
prs = Presentation()
# choose what kind of slide we would like to use 
title_slide_layout = prs.slide_layouts[5]

# I created a function to develop each and every slide for the powerpoint 

def slide(Name, date, vistor,state, coordinates, acres, km, description,index):

    # I used the slide title box to insert the Name of the park
    # I didnt need to move the box or anything   
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = Name

    # manage the height and locaiton of each element being addded 
    height = 1.8
    left = 6

    # create a textbox so we can add text 
    txBox = slide.shapes.add_textbox(Inches(left), Inches(height),Inches(1), Inches(1)) 
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    # add text to our textbox and input information about date etc, repeat for every info box 
    p.text = "Date Founded: " + date 
    
    # add space inbetween text to create 
    space =.4

    txBox = slide.shapes.add_textbox(Inches(left), Inches(height + (1*space)),Inches(1), Inches(1)) 
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Number of Vistors: " + str(vistor)

    txBox = slide.shapes.add_textbox(Inches(left), Inches(height + (2*space)),Inches(1), Inches(1)) 
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "State: " + state
    
    txBox = slide.shapes.add_textbox(Inches(left), Inches(height + (3*space)),Inches(1), Inches(1)) 
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Coordinates: " + coordinates 

    txBox = slide.shapes.add_textbox(Inches(left), Inches(height + (4*space)),Inches(1), Inches(1)) 
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Acres: " + str(acres) 

    txBox = slide.shapes.add_textbox(Inches(left), Inches(height + (5*space)),Inches(1), Inches(1)) 
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "KM Squared: " + str(km) 



    # the spacing is diffrent for this box because this is where I inputed the description of the park into the pptx 

    txBox = slide.shapes.add_textbox(Inches(1), Inches(5),Inches(8.5), Inches(4)) 
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "     " + description 

    # Taking our photo's url's we can download them as an image and then insert them into the powerpoint 
    # we have to download the images because we are unable to refrence the image just using the url  

    print(index)
    #UPDATE where you would like to save a photo 
    urllib.request.urlretrieve(index, "E:/Code Project/Python - PPTX project/park.jpg")
    
    # we will adjust the shape of the photo 
    #UPDATE WHERE YOU WOULD LIKE TO SAVE PHOTO
    p = slide.shapes.add_picture('E:/Code Project/Python - PPTX project/park.jpg', Inches(1.2), Inches(1.8),Inches(4.5),Inches(3))

# go through our list of national parks and add the necessary info 
for i in range(len(df)): 
    slide(df.iloc[i,0],df.iloc[i,1],df.iloc[i,2],df.iloc[i,5], df.iloc[i,4], df.iloc[i,6],df.iloc[i,7], df.iloc[i,3],df2.iloc[i])


# save the powerpoint in my folder  
#UPDATE 
# change to where you would like to save your PowerPoint 
prs.save('E:/Code Project/Python - PPTX project/test.pptx')

# open the powerpoint application on my computer 
import os
#UPDATE
# change to where you put your PowerPoint 
os.startfile('E:/Code Project/Python - PPTX project/test.pptx')

